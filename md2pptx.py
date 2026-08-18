"""
Convertitore Markdown -> PowerPoint (senza formattazione)
================================================================================
Questo script:
1. Legge un file .md strutturato secondo le regole fornite.
2. Genera un file .pptx con la gerarchia:
   - #  → Slide titolo
   - ## → Slide sezione (con eventuale sottotitolo)
   - ### → Slide normale (con corpo composto da testo, elenchi puntati/numerati, grassetto)

Requisiti:
    pip install python-pptx
================================================================================
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def parse_markdown(md_path):
    """
    Legge il file Markdown e restituisce una lista di dizionari, ognuno dei quali
    rappresenta una slide.
    
    Ogni slide ha:
        - type: 'title' | 'section' | 'content'
        - title: str
        - subtitle: str (solo per type=='section')
        - body: list of str (solo per type=='content')
                 Le stringhe possono essere:
                 - testo semplice
                 - testo con **grassetto** (verrà gestito in fase di generazione)
                 - elementi di elenco che iniziano con '- ', '* ' o '1. ', '2. ', ...
    """
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    slides = []
    current_h1 = None
    current_h2 = None
    current_h3 = None
    body_lines = []
    section_subtitle = None
    inside_section_description = False

    i = 0
    while i < len(lines):
        line = lines[i].rstrip('\n')

        # --- Rilevamento intestazioni ---
        h1_match = re.match(r'^#\s+(.+)$', line)
        h2_match = re.match(r'^##\s+(.+)$', line)
        h3_match = re.match(r'^###\s+(.+)$', line)

        if h1_match:
            # Chiudi eventuale slide aperta
            if current_h3 is not None:
                slides.append({'type': 'content', 'title': current_h3, 'body': body_lines})
                current_h3 = None
                body_lines = []
            if current_h2 is not None:
                # Se c'era una sezione aperta con sottotitolo, la aggiungiamo
                if section_subtitle is not None:
                    slides.append({'type': 'section', 'title': current_h2, 'subtitle': section_subtitle})
                    section_subtitle = None
                else:
                    slides.append({'type': 'section', 'title': current_h2, 'subtitle': None})
                current_h2 = None

            current_h1 = h1_match.group(1)
            slides.append({'type': 'title', 'title': current_h1})
            inside_section_description = False
            i += 1
            continue

        if h2_match:
            # Chiudi eventuale slide H3 aperta
            if current_h3 is not None:
                slides.append({'type': 'content', 'title': current_h3, 'body': body_lines})
                current_h3 = None
                body_lines = []
            # Chiudi eventuale sezione H2 aperta
            if current_h2 is not None:
                if section_subtitle is not None:
                    slides.append({'type': 'section', 'title': current_h2, 'subtitle': section_subtitle})
                    section_subtitle = None
                else:
                    slides.append({'type': 'section', 'title': current_h2, 'subtitle': None})
                current_h2 = None

            current_h2 = h2_match.group(1)
            inside_section_description = True  # il prossimo paragrafo è il sottotitolo
            i += 1
            continue

        if h3_match:
            # Chiudi eventuale slide H3 aperta
            if current_h3 is not None:
                slides.append({'type': 'content', 'title': current_h3, 'body': body_lines})
                current_h3 = None
                body_lines = []
            # Se non c'è una sezione H2 aperta, la creiamo fittizia (ma secondo le regole
            # non dovrebbe succedere, lo gestiamo comunque)
            if current_h2 is None:
                current_h2 = "Sezione"  # caso eccezionale
                # Se c'era un sottotitolo pendente dalla sezione (non dovrebbe), lo scartiamo
                section_subtitle = None

            current_h3 = h3_match.group(1)
            inside_section_description = False
            i += 1
            continue

        # --- Gestione del contenuto ---
        if line.strip() == '':
            # Riga vuota: la ignoriamo, ma se siamo in una sezione e non c'è ancora sottotitolo,
            # non facciamo nulla.
            i += 1
            continue

        # Se siamo in una sezione H2 e stiamo raccogliendo il sottotitolo
        if inside_section_description and current_h2 is not None and current_h3 is None:
            # Il sottotitolo è il primo paragrafo dopo l'H2
            section_subtitle = line.strip()
            inside_section_description = False
            i += 1
            continue

        # Altrimenti, se siamo all'interno di una slide H3, accumuliamo il corpo
        if current_h3 is not None:
            body_lines.append(line)
            i += 1
            continue

        # Se arriviamo qui, siamo in uno stato non gestito (ad esempio testo prima del primo H1)
        # lo ignoriamo
        i += 1

    # --- Chiusura ultime slide aperte ---
    if current_h3 is not None:
        slides.append({'type': 'content', 'title': current_h3, 'body': body_lines})
    if current_h2 is not None:
        if section_subtitle is not None:
            slides.append({'type': 'section', 'title': current_h2, 'subtitle': section_subtitle})
        else:
            slides.append({'type': 'section', 'title': current_h2, 'subtitle': None})

    return slides


def apply_bold_to_text(paragraph, text):
    """
    Aggiunge un paragrafo a un elemento (shape.text_frame) interpretando il
    **grassetto** Markdown. 
    Viene creato un unico paragrafo e vengono applicati i run con grassetto.
    """
    # Dividiamo il testo in base a ** ... **
    parts = re.split(r'(\*\*[^*]+\*\*)', text)
    for part in parts:
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            # Grassetto
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.bold = True
        else:
            # Testo normale
            run = paragraph.add_run()
            run.text = part
            run.bold = False


def generate_pptx(slides, output_path):
    """Crea il file PowerPoint a partire dalla lista di slide."""
    prs = Presentation()

    for slide_data in slides:
        slide_type = slide_data['type']
        title = slide_data['title']

        if slide_type == 'title':
            # Layout titolo (0 = titolo + sottotitolo)
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            # Non impostiamo il sottotitolo

        elif slide_type == 'section':
            # Layout titolo + sottotitolo
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            subtitle = slide_data.get('subtitle')
            if subtitle:
                # Il placeholder per il sottotitolo è il secondo nella slide layout 0
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = subtitle

        elif slide_type == 'content':
            # Layout titolo + contenuto (1 = titolo + contenuto)
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title

            # Il placeholder per il contenuto è il secondo nella slide layout 1
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Rimuove il placeholder predefinito

            body = slide_data.get('body', [])
            if not body:
                continue

            # Variabili per gestire elenchi
            in_bullet_list = False
            in_numbered_list = False
            list_counter = 0

            for line in body:
                line = line.strip()
                if not line:
                    continue

                # Rileviamo se è un elemento di elenco puntato
                bullet_match = re.match(r'^-\s+(.+)$', line) or re.match(r'^\*\s+(.+)$', line)
                numbered_match = re.match(r'^(\d+)\.\s+(.+)$', line)

                if bullet_match:
                    # Elenco puntato
                    text_content = bullet_match.group(1)
                    if not in_bullet_list:
                        # Se eravamo in un elenco numerato, lo chiudiamo
                        in_numbered_list = False
                        list_counter = 0
                    in_bullet_list = True
                    p = text_frame.add_paragraph()
                    p.level = 0
                    p.bullet = True
                    apply_bold_to_text(p, text_content)
                    continue

                if numbered_match:
                    # Elenco numerato
                    text_content = numbered_match.group(2)
                    if not in_numbered_list:
                        # Se eravamo in un elenco puntato, lo chiudiamo
                        in_bullet_list = False
                    in_numbered_list = True
                    # La numerazione viene gestita automaticamente da PowerPoint
                    p = text_frame.add_paragraph()
                    p.level = 0
                    p.bullet = True
                    # PowerPoint usa automaticamente i numeri se impostiamo il tipo di elenco
                    # Ma per semplicità usiamo il testo numerato manuale
                    # Oppure possiamo impostare il livello e usare il bullet automatico.
                    # Usiamo un bullet personalizzato.
                    p.text = f"{list_counter+1}. {text_content}"
                    p.bullet = False  # Disabilitiamo il bullet perché abbiamo già il numero
                    list_counter += 1
                    continue

                # Testo semplice (non elenco)
                if in_bullet_list or in_numbered_list:
                    # Se eravamo in un elenco e ora c'è testo semplice, usciamo dall'elenco
                    in_bullet_list = False
                    in_numbered_list = False
                    list_counter = 0

                p = text_frame.add_paragraph()
                p.bullet = False
                apply_bold_to_text(p, line)

    prs.save(output_path)
    print(f"✅ Presentazione salvata in: {output_path}")


def main():
    if len(sys.argv) < 2:
        print("Utilizzo: python script.py input.md [output.pptx]")
        print("  Se output.pptx non viene specificato, viene usato 'presentazione.pptx'")
        sys.exit(1)

    md_file = sys.argv[1]
    if len(sys.argv) >= 3:
        pptx_file = sys.argv[2]
    else:
        pptx_file = "presentazione.pptx"

    if not Path(md_file).exists():
        print(f"❌ File {md_file} non trovato.")
        sys.exit(1)

    print(f"📖 Lettura di {md_file}...")
    slides = parse_markdown(md_file)

    print(f"📊 Trovate {len(slides)} slide:")
    for s in slides:
        print(f"  - [{s['type']}] {s['title'][:40]}...")

    generate_pptx(slides, pptx_file)


if __name__ == "__main__":
    main()