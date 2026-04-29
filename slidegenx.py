import argparse
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title_text):
    """Aggiunge una slide solo con titolo"""
    slide_layout = prs.slide_layouts[0]  # Titolo e corpo
    slide = prs.slides.add_slide(slide_layout)
    
    # Imposta solo il titolo
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Cancella il contenuto del corpo se presente
    if slide.shapes[1].has_text_frame:
        slide.shapes[1].text_frame.clear()
    
    return slide

def add_content_slide(prs, title_text, content_lines):
    """
    Aggiunge una slide con titolo e corpo.
    content_lines: lista di stringhe o tuple (text, is_bold, is_bullet)
    """
    slide_layout = prs.slide_layouts[1]  # Titolo e contenuto
    slide = prs.slides.add_slide(slide_layout)
    
    # Imposta titolo
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Ottieni il box del contenuto
    content_box = slide.shapes[1]
    text_frame = content_box.text_frame
    text_frame.clear()
    
    for line in content_lines:
        p = text_frame.add_paragraph()
        
        if isinstance(line, tuple):
            text, is_bold, is_bullet = line
            p.text = text
            p.font.bold = is_bold
            p.font.size = Pt(24)
            if is_bullet:
                p.level = 0
                p.text = f"• {text}"
        else:
            # Se è una stringa semplice
            p.text = line
            p.font.size = Pt(24)
            p.font.bold = False
            # Verifica se inizia con bullet
            if line.startswith("•"):
                p.level = 0
    
    return slide

def parse_yaml_content(yaml_data):
    """Parsa il contenuto YAML e restituisce una lista di slide"""
    slides = []
    
    # Slide 1: Titolo del corso
    slides.append({
        "type": "title",
        "title": yaml_data["course"]["title"]
    })
    
    # Slide 2: Introduzione
    slides.append({
        "type": "content",
        "title": yaml_data["intro"]["title"],
        "content": yaml_data["intro"]["bullets"]
    })
    
    # Sezione Teoria
    theory = yaml_data["theory"]
    slides.append({
        "type": "title",
        "title": theory["section_title"]
    })
    
    for topic in theory["topics"]:
        content_items = []
        # Aggiungi definizione se presente
        if "definition" in topic:
            content_items.append((topic["definition"], True, False))
            content_items.append(("", False, False))
        
        # Aggiungi punti principali
        for point in topic.get("points", []):
            content_items.append((point, False, True))
        
        slides.append({
            "type": "content",
            "title": topic["title"],
            "content": content_items
        })
    
    # Sezione Pratica
    practice = yaml_data["practice"]
    slides.append({
        "type": "title",
        "title": practice["section_title"]
    })
    
    for activity in practice["activities"]:
        content_items = []
        for step in activity.get("steps", []):
            content_items.append((step, False, True))
        
        slides.append({
            "type": "content",
            "title": activity["title"],
            "content": content_items
        })
    
    # Caso di studio
    case_study = yaml_data["case_study"]
    slides.append({
        "type": "title",
        "title": case_study["title"]
    })
    
    for i, detail in enumerate(case_study["details"]):
        content_items = []
        if i == 0:  # Scenario
            content_items.append(("Scenario:", True, False))
            content_items.append(("", False, False))
            for line in detail["scenario"]:
                content_items.append((line, False, True))
        elif i == 1:  # Lezioni apprese
            content_items.append(("Lesson Learned:", True, False))
            content_items.append(("", False, False))
            for line in detail["lessons"]:
                content_items.append((line, False, True))
        
        slides.append({
            "type": "content",
            "title": f"{case_study['title']} - {detail['subtitle']}",
            "content": content_items
        })
    
    # Discussione finale
    discussion = yaml_data["discussion"]
    slides.append({
        "type": "content",
        "title": discussion["title"],
        "content": [(point, False, True) for point in discussion["points"]]
    })
    
    return slides

def generate_presentation(yaml_file, output_file):
    """Genera la presentazione PowerPoint"""
    # Carica YAML
    with open(yaml_file, 'r', encoding='utf-8') as f:
        yaml_data = yaml.safe_load(f)
    
    # Crea presentazione (nessun tema/colore predefinito)
    prs = Presentation()
    
    # Rimuovi tema predefinito (usa slide vuote con stile base)
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            if hasattr(shape, 'fill'):
                shape.fill.background()
    
    # Parsa contenuto
    slides_data = parse_yaml_content(yaml_data)
    
    # Genera slides
    for slide_data in slides_data:
        if slide_data["type"] == "title":
            add_title_slide(prs, slide_data["title"])
        else:
            add_content_slide(prs, slide_data["title"], slide_data["content"])
    
    # Salva presentazione
    prs.save(output_file)
    print(f"Presentazione generata: {output_file}")
    print(f"Totale slides: {len(prs.slides)}")

def main():
    parser = argparse.ArgumentParser(description="Genera slides PowerPoint da file YAML")
    parser.add_argument("yml_file", help="Percorso del file YAML di configurazione")
    parser.add_argument("-o", "--output", default="output.pptx", 
                        help="Percorso del file PowerPoint di output")
    
    args = parser.parse_args()
    
    if not os.path.exists(args.yml_file):
        print(f"Errore: File {args.yml_file} non trovato")
        return
    
    generate_presentation(args.yml_file, args.output)

if __name__ == "__main__":
    main()